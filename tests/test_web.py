"""Tests for the den-AI studio backend (model calls mocked)."""

from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

import denai.web.server as server
from denai.web.server import create_app

MOCK_ROAST = {
    "language": "en",
    "den_score": 3.8,
    "one_liner": "Twelve slides to say one number went up.",
    "brand_verdict": {"keep": False, "comment": "No brand to save."},
    "top_sins": ["Wall of text on slide 2"],
    "sections": [
        {"index": 1, "title": "Q3", "score": 4, "roast": "Meh.", "fix": "Lead with the number."}
    ],
    "summary": "Cut it down.",
}

MOCK_SPEC = {
    "kind": "pptx",
    "use_original_brand": False,
    "slides": [
        {"layout": "title", "title": "Q3: up 12%", "bullets": [], "notes": ""},
        {"layout": "content", "title": "Two drivers", "bullets": ["Acme", "Globex"], "notes": ""},
    ],
}


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Q3 Results"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def client(monkeypatch: pytest.MonkeyPatch) -> TestClient:
    monkeypatch.setattr(
        server, "roast_document", lambda extraction, model, language=None: MOCK_ROAST
    )
    monkeypatch.setattr(server, "fix_document", lambda extraction, roast, model: MOCK_SPEC)
    return TestClient(create_app())


def test_index_serves_ui(client: TestClient) -> None:
    res = client.get("/")
    assert res.status_code == 200
    assert "den-AI studio" in res.text


def test_roast_fix_download_flow(client: TestClient, sample_pptx: Path) -> None:
    with sample_pptx.open("rb") as fh:
        res = client.post(
            "/api/roast", files={"file": ("deck.pptx", fh)}, data={"language": "it"}
        )
    assert res.status_code == 200
    data = res.json()
    assert data["roast"]["den_score"] == 3.8
    job_id = data["job_id"]

    res = client.post(f"/api/fix/{job_id}")
    assert res.status_code == 200
    assert res.json()["use_original_brand"] is False

    res = client.get(f"/api/download/{job_id}")
    assert res.status_code == 200
    assert res.headers["content-disposition"].endswith('deck.denai-fixed.pptx"')


def test_rejects_unknown_extension(client: TestClient) -> None:
    res = client.post("/api/roast", files={"file": ("notes.txt", b"hello")})
    assert res.status_code == 400


def test_rejects_unknown_language(client: TestClient, sample_pptx: Path) -> None:
    with sample_pptx.open("rb") as fh:
        res = client.post(
            "/api/roast", files={"file": ("deck.pptx", fh)}, data={"language": "fr"}
        )
    assert res.status_code == 400


def test_fix_unknown_job(client: TestClient) -> None:
    assert client.post("/api/fix/nope").status_code == 404
    assert client.get("/api/download/nope").status_code == 404
