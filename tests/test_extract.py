"""Smoke tests for document extraction (no API calls)."""

from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation

from denai.extract import extract_document


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Q3 Results"
    slide.placeholders[1].text = "A subtitle nobody reads"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Revenue"
    slide2.placeholders[1].text = "Up and to the right"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def sample_docx(tmp_path: Path) -> Path:
    doc = Document()
    doc.add_heading("Monthly Report", level=1)
    doc.add_paragraph("Everything is fine.")
    doc.add_heading("KPIs", level=2)
    doc.add_paragraph("Numbers go up.")
    path = tmp_path / "report.docx"
    doc.save(str(path))
    return path


def test_extract_pptx(sample_pptx: Path) -> None:
    result = extract_document(sample_pptx)
    assert result["kind"] == "pptx"
    assert result["n_units"] == 2
    assert result["units"][0]["title"] == "Q3 Results"
    assert "brand" in result


def test_extract_docx(sample_docx: Path) -> None:
    result = extract_document(sample_docx)
    assert result["kind"] == "docx"
    assert result["n_units"] == 2
    titles = [u["title"] for u in result["units"]]
    assert "Monthly Report" in titles
    assert "KPIs" in titles


def test_unsupported_extension(tmp_path: Path) -> None:
    path = tmp_path / "file.txt"
    path.write_text("hello")
    with pytest.raises(ValueError):
        extract_document(path)
