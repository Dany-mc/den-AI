"""Tests for the surgical editor: design preserved, ops applied in place."""

from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation

from denai.editor import apply_edits
from denai.extract import extract_document


@pytest.fixture()
def deck(tmp_path: Path) -> Path:
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Q3 Update"
    s1.placeholders[1].text = "A presentation about our quarter"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Numbers"
    s2.placeholders[1].text = "A very long wall of text about synergy and alignment."
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Thank you"
    s3.placeholders[1].text = "Questions?"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


def test_edit_pptx(deck: Path, tmp_path: Path) -> None:
    out = tmp_path / "deck-fixed.pptx"
    edits = [
        {"op": "rewrite_title", "unit": 1, "text": "Q3: revenue up 12%"},
        {"op": "replace_body", "unit": 2, "bullets": ["Acme +8%", "Globex +4%"]},
        {"op": "set_notes", "unit": 2, "text": "Explain the drivers."},
        {"op": "delete_unit", "unit": 3},
        {"op": "rewrite_title", "unit": 99, "text": "ignored"},  # bad unit → skipped
    ]
    applied = apply_edits(edits, deck, out, "pptx")
    assert applied == 4

    result = extract_document(out)
    assert result["n_units"] == 2  # thank-you slide is gone
    assert result["units"][0]["title"] == "Q3: revenue up 12%"
    assert "Acme +8%" in result["units"][1]["texts"]
    assert result["units"][1]["notes"] == "Explain the drivers."


def test_edit_pptx_never_deletes_everything(deck: Path, tmp_path: Path) -> None:
    out = tmp_path / "out.pptx"
    edits = [{"op": "delete_unit", "unit": n} for n in (1, 2, 3)]
    apply_edits(edits, deck, out, "pptx")
    assert extract_document(out)["n_units"] == 3  # refused: would delete all


@pytest.fixture()
def docx_doc(tmp_path: Path) -> Path:
    doc = Document()
    doc.add_heading("Monthly Report", level=1)
    doc.add_paragraph("Everything is fine. " * 20)
    doc.add_heading("KPIs", level=2)
    doc.add_paragraph("Numbers go up.")
    doc.add_heading("Closing thoughts", level=2)
    doc.add_paragraph("Thanks for reading.")
    path = tmp_path / "report.docx"
    doc.save(str(path))
    return path


def test_edit_docx(docx_doc: Path, tmp_path: Path) -> None:
    out = tmp_path / "report-fixed.docx"
    edits = [
        {"op": "rewrite_title", "unit": 1, "text": "March: churn doubled"},
        {"op": "replace_body", "unit": 1, "bullets": ["Churn 4% → 8%", "Driven by pricing"]},
        {"op": "delete_unit", "unit": 3},
    ]
    applied = apply_edits(edits, docx_doc, out, "docx")
    assert applied == 3

    result = extract_document(out)
    titles = [u["title"] for u in result["units"]]
    assert titles == ["March: churn doubled", "KPIs"]
    assert "Churn 4% → 8%" in result["units"][0]["texts"]
    assert "Everything is fine" not in " ".join(result["units"][0]["texts"])
