"""Apply den-AI's edit plan to a copy of the original file.

This is the heart of the product: the user's design, theme, layouts and
images are preserved — only text changes and whole filler units go.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation
from pptx.util import Pt


def apply_edits(edits: list[dict[str, Any]], src: Path, out: Path, kind: str) -> int:
    """Apply the edit plan to a copy of `src`, save as `out`.

    Returns the number of operations actually applied (bad unit numbers and
    impossible ops are skipped, never fatal).
    """
    if kind == "pptx":
        return _edit_pptx(edits, src, out)
    if kind == "docx":
        return _edit_docx(edits, src, out)
    raise ValueError(f"Surgical editing supports pptx/docx, not {kind}")


# ---------------------------------------------------------------- pptx

def _edit_pptx(edits: list[dict[str, Any]], src: Path, out: Path) -> int:
    prs = Presentation(str(src))
    slides = list(prs.slides)
    applied = 0
    to_delete: set[int] = set()

    for edit in edits:
        try:
            unit = int(edit.get("unit", 0))
        except (TypeError, ValueError):
            continue
        if not 1 <= unit <= len(slides):
            continue
        slide = slides[unit - 1]
        op = edit.get("op")
        try:
            if op == "rewrite_title" and slide.shapes.title is not None:
                slide.shapes.title.text_frame.text = str(edit.get("text", ""))
                applied += 1
            elif op == "replace_body":
                body = _body_shape(slide)
                if body is not None:
                    _set_bullets(body, [str(b) for b in edit.get("bullets") or []])
                    applied += 1
            elif op == "set_notes":
                slide.notes_slide.notes_text_frame.text = str(edit.get("text", ""))
                applied += 1
            elif op == "delete_unit":
                to_delete.add(unit)
        except (AttributeError, KeyError, ValueError):
            continue

    if to_delete and len(to_delete) < len(slides):
        sld_id_lst = prs.slides._sldIdLst  # python-pptx has no public delete API
        sld_ids = list(sld_id_lst)
        for unit in to_delete:
            sld_id_lst.remove(sld_ids[unit - 1])
            applied += 1

    prs.save(str(out))
    return applied


def _body_shape(slide):
    """The content placeholder, or the meatiest non-title text shape."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph
    candidates = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape != slide.shapes.title
    ]
    if not candidates:
        return None
    return max(candidates, key=lambda s: len(s.text_frame.text))


def _set_bullets(shape, bullets: list[str]) -> None:
    tf = shape.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = bullet
        para.level = 0
        for run in para.runs:
            run.font.size = Pt(20)


# ---------------------------------------------------------------- docx

def _docx_units(doc) -> list[dict[str, Any]]:
    """Mirror extract.py's section mapping: heading + its body paragraphs."""
    units: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    for para in doc.paragraphs:
        style = (para.style.name or "").lower() if para.style else ""
        is_heading = (style.startswith("heading") or style == "title") and para.text.strip()
        if is_heading:
            if current is not None:
                units.append(current)
            current = {"heading": para, "body": []}
        else:
            if current is None:
                if not para.text.strip():
                    continue
                current = {"heading": None, "body": []}
            current["body"].append(para)
    if current is not None:
        units.append(current)
    return units


def _edit_docx(edits: list[dict[str, Any]], src: Path, out: Path) -> int:
    doc = Document(str(src))
    units = _docx_units(doc)
    applied = 0

    for edit in edits:
        try:
            unit_no = int(edit.get("unit", 0))
        except (TypeError, ValueError):
            continue
        if not 1 <= unit_no <= len(units):
            continue
        unit = units[unit_no - 1]
        op = edit.get("op")
        try:
            if op == "rewrite_title" and unit["heading"] is not None:
                unit["heading"].text = str(edit.get("text", ""))
                applied += 1
            elif op == "replace_body":
                bullets = [str(b) for b in edit.get("bullets") or []]
                anchor = unit["heading"] or (unit["body"][0] if unit["body"] else None)
                if anchor is None:
                    continue
                new_paras = []
                for bullet in bullets:
                    np = doc.add_paragraph(bullet, style="List Bullet")
                    new_paras.append(np)
                if unit["heading"] is not None:
                    ref = unit["heading"]._p
                    for np in new_paras:
                        ref.addnext(np._p)
                        ref = np._p
                else:
                    first = unit["body"][0]._p
                    for np in new_paras:
                        first.addprevious(np._p)
                for para in unit["body"]:
                    para._p.getparent().remove(para._p)
                unit["body"] = new_paras
                applied += 1
            elif op == "delete_unit" and len(units) > 1:
                for para in unit["body"]:
                    para._p.getparent().remove(para._p)
                if unit["heading"] is not None:
                    unit["heading"]._p.getparent().remove(unit["heading"]._p)
                applied += 1
            # set_notes has no docx equivalent — silently skipped
        except (AttributeError, KeyError, ValueError):
            continue

    doc.save(str(out))
    return applied
