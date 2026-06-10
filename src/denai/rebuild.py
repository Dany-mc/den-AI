"""Rebuild documents from a den-AI fix spec, applying brand or den style."""

from __future__ import annotations

import tempfile
from pathlib import Path
from typing import Any

from docx import Document
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxRGB
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# den style: the opinionated fallback identity when the original brand is beyond saving
DEN_STYLE = {
    "palette": ["#1D1B16", "#E05D44", "#8D8775"],
    "fonts": ["Georgia"],
}


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))  # type: ignore[return-value]


def _pick_brand(spec: dict[str, Any], extraction: dict[str, Any]) -> dict[str, Any]:
    original = extraction.get("brand", {})
    if spec.get("use_original_brand") and original.get("palette"):
        return {
            "palette": original["palette"],
            "fonts": original.get("fonts") or DEN_STYLE["fonts"],
        }
    return DEN_STYLE


def fixed_suffix(original_suffix: str) -> str:
    """PDF and Markdown inputs rebuild as Word documents."""
    return original_suffix if original_suffix in (".pptx", ".docx") else ".docx"


def rebuild_document(spec: dict[str, Any], extraction: dict[str, Any], out_path: Path) -> Path:
    source_kind = extraction.get("kind")
    kind = spec.get("kind") or ("docx" if source_kind in ("pdf", "md") else source_kind)
    if kind == "pptx":
        return _rebuild_pptx(spec, extraction, out_path)
    if kind == "docx":
        return _rebuild_docx(spec, extraction, out_path)
    raise ValueError(f"Unsupported rebuild kind: {kind}")


def render_chart(chart: dict[str, Any], accent_hex: str, out_path: Path) -> Path | None:
    """Render a chart spec to PNG with matplotlib. Returns None on bad specs."""
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    labels = chart.get("labels") or []
    try:
        values = [float(v) for v in (chart.get("values") or [])]
    except (TypeError, ValueError):
        return None
    if not labels or not values or len(labels) != len(values):
        return None

    fig, ax = plt.subplots(figsize=(8, 4.5), dpi=150)
    kind = chart.get("type", "bar")
    if kind == "pie":
        shades = [accent_hex] + ["#b9b3a4", "#8d8775", "#d8d3c4", "#6e6a5c"] * 3
        ax.pie(
            values,
            labels=labels,
            colors=shades[: len(values)],
            autopct="%1.0f%%",
            startangle=90,
            wedgeprops={"linewidth": 1.5, "edgecolor": "white"},
        )
    elif kind == "line":
        ax.plot(labels, values, color=accent_hex, linewidth=2.5, marker="o")
    else:
        ax.bar(labels, values, color=accent_hex, width=0.62)

    if kind != "pie":
        for spine in ("top", "right", "left"):
            ax.spines[spine].set_visible(False)
        ax.grid(axis="y", color="#00000018", linewidth=0.8)
        ax.set_axisbelow(True)
        ax.tick_params(length=0)
    title = chart.get("title")
    if title:
        ax.set_title(title, fontsize=13, fontweight="bold", loc="left", pad=14)
    fig.tight_layout()
    fig.savefig(out_path, facecolor="white", bbox_inches="tight")
    plt.close(fig)
    return out_path


def _rebuild_pptx(spec: dict[str, Any], extraction: dict[str, Any], out_path: Path) -> Path:
    brand = _pick_brand(spec, extraction)
    accent = _hex_to_rgb(brand["palette"][0])
    font = brand["fonts"][0]

    prs = Presentation()
    layouts = {
        "title": prs.slide_layouts[0],
        "content": prs.slide_layouts[1],
        "section": prs.slide_layouts[2],
    }

    for slide_spec in spec.get("slides", []):
        layout = layouts.get(slide_spec.get("layout", "content"), layouts["content"])
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title is not None:
            title_shape = slide.shapes.title
            title_shape.text = slide_spec.get("title", "")
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = font
                    run.font.color.rgb = RGBColor(*accent)
                    run.font.bold = True

        chart = slide_spec.get("chart")
        bullets = slide_spec.get("bullets") or []
        if chart:
            bullets = bullets[:2]  # leave room for the chart

        body = next(
            (ph for ph in slide.placeholders if ph.placeholder_format.idx == 1), None
        )
        if body is not None and bullets:
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = bullet
                para.level = 0
                for run in para.runs:
                    run.font.name = font
                    run.font.size = Pt(20)

        if chart:
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                png = render_chart(chart, brand["palette"][0], Path(tmp.name))
            if png is not None:
                top = Inches(3.0) if bullets else Inches(2.0)
                slide.shapes.add_picture(str(png), Inches(1.5), top, width=Inches(7))
                png.unlink(missing_ok=True)

        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(out_path))
    return out_path


def _rebuild_docx(spec: dict[str, Any], extraction: dict[str, Any], out_path: Path) -> Path:
    brand = _pick_brand(spec, extraction)
    accent = _hex_to_rgb(brand["palette"][0])
    font = brand["fonts"][0]

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = font
    style.font.size = DocxPt(11)

    title = spec.get("title")
    if title:
        heading = doc.add_heading(title, level=0)
        for run in heading.runs:
            run.font.color.rgb = DocxRGB(*accent)

    for block in spec.get("blocks", []):
        btype = block.get("type")
        if btype == "heading":
            level = min(max(int(block.get("level", 1)), 1), 3)
            heading = doc.add_heading(block.get("text", ""), level=level)
            for run in heading.runs:
                run.font.color.rgb = DocxRGB(*accent)
        elif btype == "paragraph":
            doc.add_paragraph(block.get("text", ""))
        elif btype == "bullets":
            for item in block.get("items", []):
                doc.add_paragraph(item, style="List Bullet")
        elif btype == "chart" and block.get("chart"):
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                png = render_chart(block["chart"], brand["palette"][0], Path(tmp.name))
            if png is not None:
                doc.add_picture(str(png), width=DocxInches(6))
                png.unlink(missing_ok=True)

    doc.save(str(out_path))
    return out_path
