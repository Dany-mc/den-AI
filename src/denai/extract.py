"""Extract structured content and brand signals from PPTX and DOCX files."""

from __future__ import annotations

from collections import Counter
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Colors that carry no brand signal
_NEUTRAL_COLORS = {"000000", "FFFFFF", "FEFEFE", "010101"}


SUPPORTED_SUFFIXES = (".pptx", ".docx", ".pdf", ".md")


def extract_document(path: Path) -> dict[str, Any]:
    """Return a structured view of the document: content units + brand signals."""
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".md":
        return _extract_md(path)
    raise ValueError(
        f"Unsupported file type: {suffix} (expected one of {', '.join(SUPPORTED_SUFFIXES)})"
    )


def _distill_brand(fonts: Counter, colors: Counter, has_images: bool) -> dict[str, Any]:
    """Keep only the strong, reliable signals: dominant colors and main fonts."""
    palette = [f"#{c}" for c, _ in colors.most_common(8) if c not in _NEUTRAL_COLORS][:3]
    return {
        "palette": palette,
        "fonts": [f for f, _ in fonts.most_common(3)],
        "has_images": has_images,
    }


def _extract_pptx(path: Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    fonts: Counter = Counter()
    colors: Counter = Counter()
    total_images = 0
    units: list[dict[str, Any]] = []

    for index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        n_images = 0
        n_shapes = 0
        title = ""
        try:
            if slide.shapes.title is not None:
                title = slide.shapes.title.text.strip()
        except (AttributeError, KeyError):
            pass

        for shape in slide.shapes:
            n_shapes += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                n_images += 1
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        texts.append(run.text.strip())
                    if run.font.name:
                        fonts[run.font.name] += 1
                    rgb = _run_rgb(run)
                    if rgb:
                        colors[rgb] += 1

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        total_images += n_images
        units.append(
            {
                "index": index,
                "title": title,
                "texts": texts,
                "notes": notes,
                "n_shapes": n_shapes,
                "n_images": n_images,
            }
        )

    return {
        "kind": "pptx",
        "name": path.name,
        "n_units": len(units),
        "units": units,
        "brand": _distill_brand(fonts, colors, total_images > 0),
    }


def _run_rgb(run) -> str | None:
    try:
        color = run.font.color
        if color is not None and color.type is not None and color.rgb is not None:
            return str(color.rgb)
    except (AttributeError, TypeError):
        pass
    return None


def _extract_docx(path: Path) -> dict[str, Any]:
    doc = Document(str(path))
    fonts: Counter = Counter()
    colors: Counter = Counter()
    units: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    section_index = 0

    def _new_section(title: str) -> dict[str, Any]:
        nonlocal section_index
        section_index += 1
        return {"index": section_index, "title": title, "texts": [], "notes": ""}

    for para in doc.paragraphs:
        text = para.text.strip()
        style = (para.style.name or "").lower() if para.style else ""
        for run in para.runs:
            if run.font.name:
                fonts[run.font.name] += 1
            try:
                if run.font.color is not None and run.font.color.rgb is not None:
                    colors[str(run.font.color.rgb)] += 1
            except (AttributeError, TypeError):
                pass
        if not text:
            continue
        if style.startswith("heading") or style == "title":
            if current is not None:
                units.append(current)
            current = _new_section(text)
        else:
            if current is None:
                current = _new_section("(intro)")
            current["texts"].append(text)

    if current is not None:
        units.append(current)

    n_tables = len(doc.tables)
    n_images = len(doc.inline_shapes)
    return {
        "kind": "docx",
        "name": path.name,
        "n_units": len(units),
        "n_tables": n_tables,
        "units": units,
        "brand": _distill_brand(fonts, colors, n_images > 0),
    }


def _extract_pdf(path: Path) -> dict[str, Any]:
    """One unit per page. PDFs carry no reliable brand signals — none extracted."""
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    units: list[dict[str, Any]] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        units.append(
            {
                "index": index,
                "title": lines[0][:120] if lines else f"Page {index}",
                "texts": lines[1:][:80],
                "notes": "",
            }
        )
    return {
        "kind": "pdf",
        "name": path.name,
        "n_units": len(units),
        "units": units,
        "brand": {"palette": [], "fonts": [], "has_images": False},
    }


def _extract_md(path: Path) -> dict[str, Any]:
    """Markdown headings become section boundaries."""
    units: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    in_fence = False

    def _new(title: str) -> dict[str, Any]:
        return {"index": len(units) + 1, "title": title, "texts": [], "notes": ""}

    for raw in path.read_text(encoding="utf-8", errors="replace").splitlines():
        line = raw.strip()
        if line.startswith("```"):
            in_fence = not in_fence
            continue
        if in_fence or not line:
            continue
        if line.startswith("#"):
            if current is not None:
                units.append(current)
            current = _new(line.lstrip("#").strip())
        else:
            if current is None:
                current = _new("(intro)")
            current["texts"].append(line)

    if current is not None:
        units.append(current)
    for i, unit in enumerate(units, start=1):
        unit["index"] = i

    return {
        "kind": "md",
        "name": path.name,
        "n_units": len(units),
        "units": units,
        "brand": {"palette": [], "fonts": [], "has_images": False},
    }


def extract_csv(path: Path, max_rows: int = 200) -> dict[str, Any]:
    """Tabular input for `denai report`: headers + a capped sample of rows."""
    import csv

    with path.open(newline="", encoding="utf-8", errors="replace") as fh:
        sample = fh.read(4096)
        fh.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",;\t")
        except csv.Error:
            dialect = csv.excel
        reader = csv.reader(fh, dialect)
        rows = list(reader)

    if not rows:
        raise ValueError("Empty CSV — even den-AI can't judge a void.")
    headers, body = rows[0], rows[1:]
    return {
        "kind": "csv",
        "name": path.name,
        "headers": headers,
        "n_rows": len(body),
        "rows": body[:max_rows],
    }
